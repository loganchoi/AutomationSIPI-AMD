import skrf as rf
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import os
import tkinter as tk
from tkinter import filedialog
'''
NOTE: There is a bug with the returnLoss Table where it generates the second figure listing out the each 5th frequency the highest port value.
'''


def select_directory():
    root = tk.Tk()
    root.withdraw()  # Hide the root window
    directory = filedialog.askdirectory(title="Select Directory Containing Touchstone Files")
    return directory

def plot_insertion_Check_DQ(frequency, network, axes,file,portNames):
    num_ports = network.number_of_ports

    # Assuming frequency is an array of frequencies in Hz
    min_frequency = np.min(frequency) / 1e9  # Convert to GHz
    max_frequency = np.max(frequency) / 1e9  # Convert to GHz

    # Calculate the number of 5 GHz intervals within the range
    iterables = int((max_frequency - min_frequency) / 5)

    data = [float('inf')] * iterables
    portInfo = [""] * iterables

    ranges = []
    # Plot magnitude for each port
    for port in range(num_ports):
        if port % 2 == 0:
            if "DQS" in portNames[port]:
                magnitude_data = 20 * np.log10(np.abs(network.s[:, port, port+1]))
                axes.plot(frequency / 1e9, magnitude_data, label=f'IL')
                curRange = min(magnitude_data)
                mag_index = np.where(curRange == magnitude_data)
                freqIndex = frequency[mag_index] / 1e9
                curRange = str("{:.2f}".format(curRange))
                ranges.append(curRange + " : {}GHz ".format(freqIndex[0]) + "({},{})".format(portNames[port],portNames[port+1]))
                for i in range(1, iterables + 1):
                    closest_index = np.argmin(np.abs(frequency - (i * 5) * 1e9))
                    magnitude_at_freq = magnitude_data[closest_index]
                    if data[i-1] > magnitude_at_freq:
                        data[i-1] = magnitude_at_freq
                        portInfo[i-1] = "({},{})".format(portNames[port],portNames[port+1])


    ranges.sort(key=lambda x: float(x.split(':')[0]))   

    axes.set_title(file)
    axes.set_xlabel('Frequency (GHz)')
    axes.set_ylabel('Magnitude (dB)')
    axes.grid(True)

    # Annotate the first five ranges on the bottom left corner
    for i, r in enumerate(ranges[:3]):
        xVal, yVal = 0,0
        yVal = float(ranges[i].split(" ")[0])
        xVal = float(ranges[i].split(" ")[2].replace("GHz",""))

         # Set the offset for the annotation text
        x_offset = 0.01  # Adjust as needed
        y_offset = i * 0.15  # Adjust as needed for vertical stacking

        axes.annotate(r, xy=(xVal, yVal), xytext=(x_offset, y_offset),
                    arrowprops=dict(facecolor='black', arrowstyle='->', connectionstyle="angle,angleA=0,angleB=90,rad=0"),
                    xycoords='data', textcoords='axes fraction', ha='left', va='bottom')

    
    return data,portInfo

def plot_return(frequency, network, axes,file,portNames):
    num_ports = network.number_of_ports

    # Assuming frequency is an array of frequencies in Hz
    min_frequency = np.min(frequency) / 1e9  # Convert to GHz
    max_frequency = np.max(frequency) / 1e9  # Convert to GHz

    # Calculate the number of 5 GHz intervals within the range
    iterables = int((max_frequency - min_frequency) / 5)

    data = [float('-inf')] * iterables
    portInfo = [""] * iterables

    maximums = []
    # Plot magnitude for each port and measure Return Loss and find the worst 3
    for port in range(num_ports):
        if "DQS" in portNames[port]: #Try to implement user input for DQS or Checks or DQ
            magnitude_data = 20 * np.log10(np.abs(network.s[:, port, port]))
            axes.plot(frequency / 1e9, magnitude_data, label=f'RL')
            curMax = max(magnitude_data)  # Find the maximum excluding 0-5 GHz
            mag_index = np.where(curMax == magnitude_data)
            freqIndex = frequency[mag_index] / 1e9
            curMax = str("{:.2f}".format(curMax))
            maximums.append(curMax + " : {}GHz ".format(freqIndex[0]) + "({})".format(portNames[port]))
            for i in range(1, iterables + 1):
                    closest_index = np.argmax(np.abs(frequency - (i * 5) * 1e9))
                    magnitude_at_freq = magnitude_data[closest_index]
                    if data[i-1] < magnitude_at_freq:
                        data[i-1] = magnitude_at_freq
                        portInfo[i-1] = "({})".format(portNames[port])
    
    maximums.sort(key=lambda x: float(x.split(':')[0]),reverse=True)   

    axes.set_title(file)
    axes.set_xlabel('Frequency (GHz)')
    axes.set_ylabel('Magnitude (dB)')
    axes.grid(True)

    # Annotate the first five ranges on the bottom left corner
    for i, r in enumerate(maximums[:3]):
        xVal, yVal = 0,0
        yVal = float(maximums[i].split(" ")[0])
        xVal = float(maximums[i].split(" ")[2].replace("GHz",""))

         # Set the offset for the annotation text
        x_offset = 1  # Adjust as needed
        y_offset = i * 0.15  # Adjust as needed for vertical stacking

        axes.annotate(r, xy=(xVal, yVal), xytext=(x_offset, y_offset),
                    arrowprops=dict(facecolor='black', arrowstyle='->', connectionstyle="arc3,rad=0"),
                    xycoords='data', textcoords='axes fraction', ha='right', va='bottom')

    
    return data,portInfo


if __name__ == "__main__":
    # Input the directory containing the touchstone files
    directory = select_directory()
    file_paths = [os.path.join(directory, filename) for filename in os.listdir(directory)]

    # Create a PowerPoint presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Use the layout suitable for content

    # Create a subplot for each magnitude plot
    fig, axs = plt.subplots(4, 2, figsize=(12, 8))

    allData = []
    allPortInfo = []

    for i, file_path in enumerate(file_paths):
        file = file_path.split("\\")[-1]
        # Read network data from touchstone file
        network = rf.Network(file_path)
        portNames = network.port_names
        frequency = network.f

        # Plot insertion loss on the corresponding axes
        plot_row = i // 2
        plot_col = i % 2
        data, portInfo = plot_insertion_Check_DQ(frequency, network, axs[plot_row, plot_col],file,portNames)
        allData.append(data)
        allPortInfo.append(portInfo)

    # Adjust layout and save the plot as an image
    plt.tight_layout()
    plt.savefig('insertion_plot.png')

    # Add a slide for insertion loss to the presentation
    slide_insertion = prs.slides.add_slide(slide_layout)
    title_shape_insertion = slide_insertion.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_shape_insertion.text_frame.text = "Insertion Loss"
    slide_insertion.shapes.add_picture('insertion_plot.png', Inches(1), Inches(1))

    #Create a new slide layout for insertion loss table results
    slide_layout_insertion_table = prs.slide_layouts[6]

    #IMPLEMENT HERE
    rows = len(allData[0])

    plt.figure(figsize=(12,6))
    
    for d in range(len(allData)):
        # Create a subplot
        plt.subplot(4, 2, d+1)
        plt.axis('off')  # Turn off axis
        
        # Create a text string to display insertion loss data
        text = file_paths[d].split("\\")[-1] +"\n"
        for x in range(rows):
            text += "{} GHz: {:.2f}, {}\n".format((x+1) * 5, allData[d][x], allPortInfo[d][x])
        
        # Add the text to the subplot
        plt.text(0.5, 0.5, text, horizontalalignment='center', verticalalignment='center', fontsize=10)

    # Adjust layout to prevent overlap
    plt.tight_layout()

    # Save the plot as PNG
    plt.savefig('insertionLossTable.png')
    plt.show()
    
    # Add a slide for insertion loss results to the presentation
    slide_insertion = prs.slides.add_slide(slide_layout)
    title_shape_insertion = slide_insertion.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_shape_insertion.text_frame.text = "Insertion Loss Data"
    slide_insertion.shapes.add_picture('insertionLossTable.png', Inches(1), Inches(1))
    


    # Create a new slide layout for return loss
    slide_layout_return = prs.slide_layouts[6]

    # Create a subplot for each magnitude plot for return loss
    fig_return, axs_return = plt.subplots(4, 2, figsize=(12, 8))

    allData = []
    allPortInfo = []

    for i, file_path in enumerate(file_paths):
        # Read network data from touchstone file
        network = rf.Network(file_path)
        frequency = network.f

        # Plot return loss on the corresponding axes
        plot_row = i // 2
        plot_col = i % 2
        data, portInfo = plot_return(frequency, network, axs_return[plot_row, plot_col],file,portNames)

        allData.append(data)
        allPortInfo.append(portInfo)

    # Adjust layout and save the plot as an image for return loss
    plt.tight_layout()
    plt.savefig('return_plot.png')

    # Add a slide for return loss to the presentation
    slide_return = prs.slides.add_slide(slide_layout_return)
    title_shape_return = slide_return.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_shape_return.text_frame.text = "Return Loss"
    slide_return.shapes.add_picture('return_plot.png', Inches(1), Inches(1))


    slide_layout_return_table = prs.slide_layouts[6]

    plt.figure(figsize=(12,6))
    
    for d in range(len(allData)):
        # Create a subplot
        plt.subplot(4, 2, d+1)
        plt.axis('off')  # Turn off axis
        
        # Create a text string to display insertion loss data
        text = file_paths[d].split("\\")[-1] +"\n"
        for x in range(rows):
            text += "{} GHz: {:.2f}, {}\n".format((x+1) * 5, allData[d][x], allPortInfo[d][x])
        
        # Add the text to the subplot
        plt.text(0.5, 0.5, text, horizontalalignment='center', verticalalignment='center', fontsize=10)

    # Adjust layout to prevent overlap
    plt.tight_layout()

    # Save the plot as PNG
    plt.savefig('returnLossTable.png')
    plt.show()
    
    slide_insertion = prs.slides.add_slide(slide_layout)
    title_shape_insertion = slide_insertion.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_shape_insertion.text_frame.text = "Return Loss Data"
    slide_insertion.shapes.add_picture('returnLossTable.png', Inches(1), Inches(1))
    

    # Save the PowerPoint presentation
    prs.save('test.pptx')

    plt.close()
